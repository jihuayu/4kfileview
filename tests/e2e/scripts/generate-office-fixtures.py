#!/usr/bin/env python3
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

fixtures = Path(__file__).resolve().parent.parent / "fixtures"
fixtures.mkdir(parents=True, exist_ok=True)

# DOCX
_doc = Document()
_doc.add_heading("kkFileView E2E", level=1)
_doc.add_paragraph("This is a DOCX fixture for Phase-2 E2E.")
_doc.save(fixtures / "sample.docx")

# XLSX
_wb = Workbook()
_ws = _wb.active
_ws.title = "Sheet1"
_ws["A1"] = "name"
_ws["B1"] = "value"
_ws["A2"] = "kkFileView"
_ws["B2"] = 2
_wb.save(fixtures / "sample.xlsx")

# PPTX
_prs = Presentation()
slide_layout = _prs.slide_layouts[1]
slide = _prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "kkFileView E2E"
slide.placeholders[1].text = "This is a PPTX fixture for Phase-2 E2E."
_prs.save(fixtures / "sample.pptx")

print("office fixtures generated in", fixtures)
